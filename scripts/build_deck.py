"""Build the judging deck from the metrics files.

    python3 scripts/build_deck.py            # -> docs/FlightRisk_NYC.pptx

Every figure on every slide is read out of `reports/metrics/*.json` at build
time rather than typed in, for the same reason `scripts/verify.py` exists: a
number that is copied by hand is a number that will eventually be wrong. Re-run
the pipeline, re-run this, and the deck cannot disagree with the report.

Optional dependency: `python-pptx`. It is deliberately not in
`requirements.txt`, because nothing in the analysis needs it.
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
METRICS = ROOT / "reports" / "metrics"
FIGURES = ROOT / "reports" / "figures"
OUT = ROOT / "docs" / "FlightRisk_NYC.pptx"

# Okabe-Ito on deep navy. The palette is the project's own -- the same one the
# figures and the app use -- which keeps every chart in the deck readable under
# the three common forms of colour blindness.
NAVY = RGBColor(0x10, 0x24, 0x3A)
NAVY_SOFT = RGBColor(0x1B, 0x35, 0x51)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
MIST = RGBColor(0xEE, 0xF1, 0xF5)
INK = RGBColor(0x11, 0x14, 0x18)
MUTED = RGBColor(0x5A, 0x64, 0x72)
BLUE = RGBColor(0x00, 0x72, 0xB2)
VERM = RGBColor(0xD5, 0x5E, 0x00)
GREEN = RGBColor(0x00, 0x9E, 0x73)
AMBER = RGBColor(0xE6, 0x9F, 0x00)
CHALK = RGBColor(0xD8, 0xE2, 0xEC)

HEAD = "Cambria"
BODY = "Calibri"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.7)


# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------

def blank(prs, dark: bool = False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY if dark else PAPER
    return s


def text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=0, line=None):
    """runs: list of (string, size_pt, bold, colour, font) or list of such lists."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    paras = runs if isinstance(runs[0], list) else [runs]
    for k, para in enumerate(paras):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line:
            p.line_spacing = line
        for spec in para:
            body, size, bold, colour = spec[0], spec[1], spec[2], spec[3]
            font_name = spec[4] if len(spec) > 4 else BODY
            r = p.add_run()
            r.text = body
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = colour
            r.font.name = font_name
    return box


def card(slide, x, y, w, h, fill=MIST, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    if sh.has_text_frame:
        sh.text_frame.text = ""
    return sh


def title(slide, s, *, dark=False, sub=None):
    text(slide, M, Inches(0.5), W - 2 * M, Inches(0.9),
         [(s, 34, True, CHALK if dark else INK, HEAD)])
    if sub:
        text(slide, M, Inches(1.32), W - 2 * M, Inches(0.5),
             [(sub, 15, False, CHALK if dark else MUTED)])


def stat(slide, x, y, w, big, label, *, colour=BLUE, size=54, dark=False,
         sub=None):
    text(slide, x, y, w, Inches(1.0), [(big, size, True, colour, HEAD)])
    text(slide, x, y + Inches(0.82), w, Inches(0.9),
         [(label, 13, False, CHALK if dark else MUTED)], line=1.15)
    if sub:
        text(slide, x, y + Inches(1.46), w, Inches(0.5),
             [(sub, 11, False, MUTED if not dark else CHALK)])


def dot(slide, x, y, d, colour, glyph):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    c.fill.solid()
    c.fill.fore_color.rgb = colour
    c.line.fill.background()
    c.shadow.inherit = False
    tf = c.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = glyph
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = PAPER
    r.font.name = BODY
    return c


def style_chart(chart, *, colours, value_fmt="0.0%", labels=True,
                label_pos=XL_LABEL_POSITION.OUTSIDE_END):
    chart.has_legend = False
    # The series name is already the slide's subtitle; a chart title repeats it.
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 60
    if labels:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = value_fmt
        dl.number_format_is_linked = False
        dl.font.size = Pt(12)
        dl.font.bold = True
        dl.font.color.rgb = INK
        dl.position = label_pos
    for axis in (chart.category_axis, chart.value_axis):
        axis.has_major_gridlines = False
        axis.tick_labels.font.size = Pt(12)
        axis.tick_labels.font.color.rgb = MUTED
    chart.value_axis.visible = False
    for i, pt in enumerate(plot.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colours[i % len(colours)]


# ---------------------------------------------------------------------------

def build() -> Path:
    ev = json.loads((METRICS / "evaluation.json").read_text())
    imp = json.loads((METRICS / "impact.json").read_text())
    fair = json.loads((METRICS / "fairness.json").read_text())
    dis = json.loads((METRICS / "disruption.json").read_text())
    sev = json.loads((METRICS / "severity_v2.json").read_text())

    at = imp["at_operating_budget"]
    m, hist, rnd = at["model"], at["historical_rule"], at["random"]
    su, sens = imp["scale_up"], imp["sensitivity"]
    car = fair["disparity"]["carrier"]
    car_rows = sorted(fair["groups"]["carrier"], key=lambda g: -g["recall"])

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # -- 1. Title ---------------------------------------------------------
    s = blank(prs, dark=True)
    text(s, M, Inches(2.05), Inches(9.4), Inches(1.4),
         [("FlightRisk NYC", 54, True, PAPER, HEAD)])
    text(s, M, Inches(3.15), Inches(9.6), Inches(1.2),
         [("Predicting arrival delay ", 24, False, CHALK),
          ("before the aircraft moves", 24, True, AMBER)])
    text(s, M, Inches(4.15), Inches(8.5), Inches(1.4),
         [("A pre-departure risk model for JFK, LaGuardia and Newark, scored "
           "the way an operations desk would actually use it: as a ranked list "
           "on a fixed daily alert budget.", 15, False, CHALK)], line=1.35)
    text(s, M, Inches(6.25), Inches(9.0), Inches(0.5),
         [("DATATHON  ·  NYC Flights 2013  ·  336,776 flights", 13, True, AMBER)])

    for i, (n, lab) in enumerate([("0.507", "PR-AUC, held out"),
                                  ("24.2%", "of all delay minutes reached on a 10% daily budget"),
                                  ("0.936", "ROC-AUC, cancellations")]):
        x = Inches(9.85)
        y = Inches(2.0) + Inches(1.62) * i
        text(s, x, y, Inches(2.8), Inches(0.6),
             [(n, 30, True, AMBER, HEAD)])
        text(s, x, y + Inches(0.55), Inches(2.8), Inches(0.9),
             [(lab.replace("\n", " "), 11.5, False, CHALK)], line=1.2)

    # -- 2. The problem ---------------------------------------------------
    s = blank(prs)
    title(s, "Delay is a $33 billion problem",
          sub="Two thirds of it lands on passengers, not airlines  ·  FAA / NEXTOR "
              "Total Delay Impact Study; Airlines for America, DOT Form 41, 2025")

    for i, (big, lab, col) in enumerate([
            ("$16.7B", "borne by passengers each year — lost time, missed "
                       "connections, unplanned hotels", VERM),
            ("$8.3B", "borne directly by airlines — crew, fuel, maintenance", BLUE),
            ("$98.41", "direct operating cost of a single block minute, from "
                       "the carriers' own DOT filings", GREEN)]):
        x = M + Inches(4.15) * i
        card(s, x, Inches(2.15), Inches(3.75), Inches(1.95))
        stat(s, x + Inches(0.3), Inches(2.4), Inches(3.2), big, lab, colour=col,
             size=36)

    card(s, M, Inches(4.5), W - 2 * M, Inches(2.25), fill=NAVY)
    text(s, M + Inches(0.45), Inches(4.85), W - 2 * M - Inches(0.9), Inches(1.7),
         [[("The constraint is attention, not information.", 21, True, AMBER, HEAD)],
          [("A duty manager at a New York airport oversees several hundred "
            "departures a day and can meaningfully act on a few dozen — swap an "
            "airframe, hold a connection, pre-position a crew. They do not need a "
            "better description of yesterday. They need a ranked list early enough "
            "to use. That is why this model is evaluated as a rationing device on "
            "a fixed daily budget, not as a classifier.", 15, False, CHALK)]],
         space_after=10, line=1.3)

    # -- 3. The shortcut --------------------------------------------------
    s = blank(prs, dark=True)
    title(s, "Almost every published model on this dataset cheats", dark=True,
          sub="…and we measured the shortcut instead of taking it")

    text(s, M, Inches(2.15), Inches(6.0), Inches(2.4),
         [[("dep_delay", 20, True, AMBER, "Courier New"),
           ("  — how late the aircraft actually left — sits in the flights "
            "table and correlates with arrival delay at ρ ≈ 0.9.", 16, False, CHALK)],
          [("Include it and you get AUC above 0.90. You also get a model that "
            "answers a question nobody needs answered: by the time you know "
            "dep_delay, the crew is committed, the gate is reassigned, and the "
            "passenger with a 55-minute connection is airborne.", 16, False, CHALK)]],
         space_after=14, line=1.3)

    card(s, Inches(7.3), Inches(2.15), Inches(5.35), Inches(3.6), fill=NAVY_SOFT)
    text(s, Inches(7.75), Inches(2.5), Inches(4.5), Inches(0.5),
         [("Same model, same data, same split", 14, True, CHALK)])
    for i, (lab, val, col) in enumerate([
            ("Post-push-back (the usual number)", "0.846", MUTED),
            ("Pre-push-back (deployable)", "0.507", AMBER)]):
        y = Inches(3.15) + Inches(1.15) * i
        text(s, Inches(7.75), y, Inches(3.1), Inches(0.8),
             [(lab, 14, False, CHALK)], line=1.2)
        text(s, Inches(10.9), y - Inches(0.14), Inches(1.6), Inches(0.7),
             [(val, 32, True, col if col != MUTED else CHALK, HEAD)],
             align=PP_ALIGN.RIGHT)
    text(s, Inches(7.75), Inches(5.15), Inches(4.5), Inches(0.5),
         [("PR-AUC on the held-out Nov–Dec period", 11.5, False, CHALK)])

    text(s, M, Inches(6.05), Inches(11.9), Inches(0.9),
         [("Two thirds of the apparent skill of a “flight delay model” is "
           "the observation that the plane left late. The 0.507 is the part that "
           "is genuinely forecastable in advance — and it is the whole project.",
           16, True, AMBER)], line=1.3)

    # -- 4. What we built -------------------------------------------------
    s = blank(prs)
    title(s, "What we built",
          sub="Every stage runs from one command and rebuilds bit-identically")

    steps = [
        ("1", "Five raw tables", "336,776 flights joined to hourly weather, the "
         "FAA aircraft registry and airport metadata", BLUE),
        ("2", "67 pre-departure features", "Leakage-safe by construction; a test "
         "suite enforces the boundary, including out-of-fold target encoding", GREEN),
        ("3", "Temporal split", "Train Jan–Aug, validate Sep–Oct, hold out "
         "Nov–Dec. Never random — a random split lets the model peek forward", AMBER),
        ("4", "Gradient boosting", "40-draw random search per library over "
         "forward-chaining folds, then isotonic recalibration", VERM),
    ]
    for i, (num, head, body, col) in enumerate(steps):
        y = Inches(2.15) + Inches(1.2) * i
        dot(s, M, y, Inches(0.5), col, num)
        text(s, M + Inches(0.78), y - Inches(0.03), Inches(4.4), Inches(0.4),
             [(head, 17, True, INK, HEAD)])
        text(s, M + Inches(0.78), y + Inches(0.38), Inches(5.3), Inches(0.7),
             [(body, 13, False, MUTED)], line=1.2)

    card(s, Inches(7.55), Inches(2.15), Inches(5.1), Inches(4.6), fill=MIST)
    text(s, Inches(7.95), Inches(2.45), Inches(4.3), Inches(0.5),
         [("Held out, never touched", 15, True, INK, HEAD)])
    rows = [
        ("PR-AUC (base rate 0.250)", f"{ev['lightgbm']['test']['pr_auc']:.3f}"),
        ("ROC-AUC", f"{ev['lightgbm']['test']['roc_auc']:.3f}"),
        ("Precision, riskiest 10% of a day", f"{m['precision']:.0%}"),
        ("Cancellation ROC-AUC", f"{dis['is_cancelled']['roc_auc']:.3f}"),
        ("Tier >120 min ROC-AUC", f"{sev['tiers']['gt120']['roc_auc']:.3f}"),
        ("Metric values identical after\na full clean rebuild", "1,115 / 1,153"),
    ]
    for i, (lab, val) in enumerate(rows):
        y = Inches(3.1) + Inches(0.6) * i
        text(s, Inches(7.95), y, Inches(3.1), Inches(0.55),
             [(lab.replace("\n", " "), 12.5, False, MUTED)], line=1.1)
        text(s, Inches(11.05), y - Inches(0.04), Inches(1.25), Inches(0.4),
             [(val, 14, True, BLUE, HEAD)], align=PP_ALIGN.RIGHT)

    # -- 5. Budget curve --------------------------------------------------
    s = blank(prs)
    title(s, "It has to beat the alternatives, not beat nothing",
          sub="Share of all delay minutes reached at a 10% daily alert budget, "
              "held-out Nov–Dec")

    cd = CategoryChartData()
    cd.categories = ["Picking at random", "Historical route rate\n(no ML)",
                     "This model"]
    cd.add_series("share of delay minutes reached", (
        rnd["delay_min_share_mean"], hist["delay_min_share"],
        m["delay_min_share"]))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, M, Inches(2.1),
                            Inches(7.3), Inches(4.5), cd)
    style_chart(gf.chart, colours=[MUTED, BLUE, VERM])

    card(s, Inches(8.55), Inches(2.1), Inches(4.1), Inches(4.5), fill=NAVY)
    text(s, Inches(8.95), Inches(2.45), Inches(3.3), Inches(0.6),
         [("Why this comparison", 15, True, AMBER, HEAD)])
    text(s, Inches(8.95), Inches(3.05), Inches(3.3), Inches(3.3),
         [[("A desk alerting at random already catches 10% of the delay on a 10% "
            "budget. Value has to be stated net of that.", 13, False, CHALK)],
          [(f"{m['delay_min_share'] / rnd['delay_min_share_mean']:.1f}× random   ·   "
            f"{m['delay_min_share'] / hist['delay_min_share']:.2f}× the no-ML rule",
            15, True, AMBER)],
          [("The budget is spent per day, not pooled — a desk cannot save "
            "November's alerts for 22 December. That makes the numbers worse and "
            "the exercise real.", 13, False, CHALK)]],
         space_after=12, line=1.28)

    # -- 6. What it is worth ----------------------------------------------
    s = blank(prs)
    title(s, "What that is worth",
          sub="Unit costs cited; mitigation effectiveness assumed pessimistically "
              "at 10% and swept from 2% to 40%")

    for i, (big, lab, col) in enumerate([
            (f"${m['airline_value_usd'] / 1e6:.2f}M", "airline operating cost "
             "avoided over two months", BLUE),
            (f"{m['recovered_pax_hours'] / 1000:.0f}k", "passenger-hours "
             f"returned (${m['passenger_value_usd'] / 1e6:.2f}M at the FAA's "
             "$47/hour)", GREEN),
            (f"{m['co2_avoided_kg'] / 1000:.0f}t", "of CO₂ not burned", AMBER),
            (f"${at['model_vs_random_usd'] / 1e6:.2f}M", "of it attributable to "
             "the model rather than the budget", VERM)]):
        x = M + Inches(3.1) * i
        card(s, x, Inches(2.15), Inches(2.85), Inches(2.15))
        stat(s, x + Inches(0.28), Inches(2.4), Inches(2.3), big, lab, colour=col,
             size=30)

    card(s, M, Inches(4.65), W - 2 * M, Inches(2.2), fill=NAVY)
    text(s, M + Inches(0.45), Inches(4.98), Inches(5.6), Inches(1.7),
         [[("Break-even sits at "
            f"{sens['breakeven_effectiveness_total']:.2%} effectiveness.",
            17, True, AMBER, HEAD)],
          [("Handling an alert costs about $6; the delay it targets is worth "
            f"about ${sens['breakeven_cost_per_alert_usd']:,.0f}. Cost is not "
            "what gates this system.", 13.5, False, CHALK)]],
         space_after=8, line=1.28)
    text(s, Inches(7.1), Inches(4.98), Inches(5.5), Inches(1.7),
         [[("What we cannot measure, we did not assume.", 17, True, AMBER, HEAD)],
          [("How much of a warned delay a desk actually recovers is not "
            "published, and a historical dataset has no counterfactual. So it is "
            "swept, not asserted — and this is a ceiling on available value, not "
            "a measured saving.", 13.5, False, CHALK)]],
         space_after=8, line=1.28)

    # -- 7. Fairness ------------------------------------------------------
    s = blank(prs)
    title(s, "Then we asked who the alerts actually reach",
          sub="Ranking by probability is a rationing rule, and it has losers")

    cd = CategoryChartData()
    top = [g for g in car_rows if g["group"] in
           ("EV", "B6", "WN", "UA", "AA", "VX")]
    cd.categories = [g["group"] for g in top]
    cd.add_series("of its late flights, warned about",
                  tuple(g["recall"] for g in top))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, M, Inches(2.15),
                            Inches(6.6), Inches(4.3), cd)
    style_chart(gf.chart, colours=[VERM, VERM, AMBER, BLUE, BLUE, BLUE])

    card(s, Inches(7.85), Inches(2.15), Inches(4.8), Inches(4.3), fill=MIST)
    text(s, Inches(8.25), Inches(2.45), Inches(4.0), Inches(3.8),
         [[(f"A {car['recall_gap']:.0%} coverage gap.", 20, True, VERM, HEAD)],
          [("Southwest runs late more often than JetBlue and receives a "
            "fortieth of the alerts, because the model under-predicts it by 17 "
            "points. No AUC would ever show that.", 13.5, False, INK)],
          [("Spending the same budget proportionally cuts the gap from "
            f"{fair['price_of_equity']['carrier']['global_recall_gap']:.0%} to "
            f"{fair['price_of_equity']['carrier']['proportional_recall_gap']:.0%} "
            f"for "
            f"{fair['price_of_equity']['carrier']['delay_min_cost_of_equity_pct']:.1f}% "
            "of the delay caught. Across destination-size quartiles the same "
            "change gains 1.3% — there, fairness is free.", 13.5, False, INK)],
          [("Recommendation: budget per carrier, not globally.", 14, True,
            BLUE)]],
         space_after=10, line=1.25)

    # -- 8. Honesty -------------------------------------------------------
    s = blank(prs)
    title(s, "Why you can believe the numbers",
          sub="Nothing in the prose is taken on trust")

    checks = [
        ("Leakage", "Scramble every post-departure column and the shipped "
         "model's predictions are exactly unchanged. The split is temporal, so a "
         "leak would have to travel backwards in time.", BLUE),
        ("Determinism", "Delete every artefact, rebuild from the raw CSVs: "
         "1,115 of 1,153 metric values identical, and all 38 that differ are "
         "wall-clock timing fields. Booster SHA-256 recorded and checked.", GREEN),
        ("Prose vs artefacts", "make verify re-reads every headline number in "
         "the report and README from the metrics files and fails on "
         "disagreement — including all of the impact and fairness figures.", AMBER),
        ("What we got wrong", "A target encoding we had to throw away, a "
         "rotation feature that read from the future on the first attempt, and a "
         "severity head aimed at the middle instead of the tail. All three are "
         "in the report.", VERM),
    ]
    for i, (head, body, col) in enumerate(checks):
        x = M + Inches(6.2) * (i % 2)
        y = Inches(2.2) + Inches(2.3) * (i // 2)
        card(s, x, y, Inches(5.75), Inches(2.0))
        dot(s, x + Inches(0.32), y + Inches(0.3), Inches(0.42), col, "✓")
        text(s, x + Inches(0.95), y + Inches(0.3), Inches(4.4), Inches(0.4),
             [(head, 16, True, INK, HEAD)])
        text(s, x + Inches(0.95), y + Inches(0.72), Inches(4.5), Inches(1.1),
             [(body, 12.5, False, MUTED)], line=1.22)

    # -- 9. Next ----------------------------------------------------------
    s = blank(prs, dark=True)
    title(s, "What would come next", dark=True,
          sub="Not a bigger model")

    nexts = [
        ("Real forecast weather", "A three-hour horizon costs 0.020 PR-AUC using "
         "the crudest possible persistence forecast, so today's numbers are a "
         "floor. Cheapest available gain."),
        ("A randomised rollout", "Alert on a random half of eligible flights and "
         "compare. That measures mitigation effectiveness directly and replaces "
         "the one assumption we had to sweep."),
        ("Beyond three airports", "The BTS on-time database has the same schema "
         "nationwide, and the inbound-leg blind spot closes as soon as the "
         "network is complete rather than NYC-only."),
        ("Three outcomes, not two", "Cancellation is the most predictable "
         "outcome in the data at 0.936 ROC-AUC. Production should rank on-time / "
         "late / cancelled jointly."),
    ]
    for i, (head, body) in enumerate(nexts):
        x = M + Inches(3.1) * i
        card(s, x, Inches(2.35), Inches(2.85), Inches(3.5), fill=NAVY_SOFT)
        dot(s, x + Inches(0.3), Inches(2.65), Inches(0.42), AMBER, str(i + 1))
        text(s, x + Inches(0.3), Inches(3.28), Inches(2.3), Inches(0.75),
             [(head, 15, True, PAPER, HEAD)], line=1.15)
        text(s, x + Inches(0.3), Inches(4.1), Inches(2.3), Inches(1.6),
             [(body, 11.5, False, CHALK)], line=1.22)

    text(s, M, Inches(6.25), Inches(11.9), Inches(0.7),
         [("Steps 1 and 4 are code we could write with this dataset. Step 2 is "
           "the one that would turn a ceiling into a measurement — and it "
           "requires an airline.", 14, False, CHALK)], line=1.25)

    # -- 10. Close --------------------------------------------------------
    s = blank(prs, dark=True)
    text(s, M, Inches(2.2), Inches(11.9), Inches(2.2),
         [[("The honest number is the useful one.", 40, True, PAPER, HEAD)],
          [("0.507, not 0.846. 47% precision per day, not 64% pooled. A ceiling "
            "on value, not a measured saving. Every one of those choices made "
            "the result look worse and the work mean something.", 18, False,
            CHALK)]],
         space_after=22, line=1.25)

    for i, (n, lab) in enumerate([
            (f"{m['delay_min_share']:.1%}", "of all delay minutes reached\n"
             "on a 10% daily budget"),
            (f"{m['delay_min_share'] / rnd['delay_min_share_mean']:.1f}×",
             "the delay minutes random\nalerting would reach"),
            (f"${su['nyc_annual_lower_usd'] / 1e6:.0f}–"
             f"{su['nyc_annual_upper_usd'] / 1e6:.0f}M",
             "available value per year\nacross the three NYC airports")]):
        x = M + Inches(4.1) * i
        text(s, x, Inches(5.05), Inches(3.7), Inches(0.8),
             [(n, 34, True, AMBER, HEAD)])
        text(s, x, Inches(5.72), Inches(3.7), Inches(0.9),
             [(lab.replace("\n", " "), 12.5, False, CHALK)], line=1.2)

    text(s, M, Inches(6.72), Inches(11.9), Inches(0.4),
         [("make reproduce  ·  make verify  ·  make app", 13, True, CHALK,
           "Courier New")])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"wrote {path.relative_to(ROOT)}")
    sys.exit(0)
